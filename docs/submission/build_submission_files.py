from __future__ import annotations

from pathlib import Path
from textwrap import dedent

import matplotlib.pyplot as plt
from PIL import Image, ImageDraw, ImageFilter, ImageOps
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Inches as PptInches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "docs" / "submission" / "generated"
ASSETS = OUT / "assets"
OUT.mkdir(parents=True, exist_ok=True)
ASSETS.mkdir(parents=True, exist_ok=True)

TITLE = "Radius"
SUBTITLE = "AI Assisted Hyperlocal Marketplace"
TEAM = "[Team Name]"
MEMBERS = "[Names and IDs]"
COURSE = "[Course Code]"
INSTRUCTOR = "[Instructor Name]"
DATE = "[Submission Date]"

COLORS = {
    "bg": (13, 10, 18),
    "surface": (28, 22, 38),
    "surface2": (39, 31, 52),
    "teal": (0, 217, 192),
    "blue": (59, 130, 255),
    "pink": (255, 61, 129),
    "amber": (255, 197, 49),
    "danger": (255, 83, 112),
    "text": (245, 243, 255),
    "muted": (166, 154, 188),
}


def rgb(name: str) -> PptRGB:
    return PptRGB(*COLORS[name])


def doc_rgb(name: str) -> RGBColor:
    return RGBColor(*COLORS[name])


def screenshot_path(*candidates: str) -> Path | None:
    for candidate in candidates:
        path = ROOT / candidate
        if path.exists():
            return path
    return None


SCREENSHOTS = {
    "hero": screenshot_path("tmp/product-pass/01-hero.png", "tmp/fix-audit/07-radar-final.png", "tmp/screenshots/1_browse_guest.png"),
    "admin": screenshot_path("tmp/product-pass/06-admin-final.png", "tmp/fix-audit/05-admin.png", "tmp/screenshots/6_admin_dashboard.png"),
    "detail": screenshot_path("tmp/fix-audit/03-detail.png", "tmp/screenshots/3_detail_page.png"),
    "sell": screenshot_path("tmp/fix-audit/04-sell.png", "tmp/screenshots/4_sell_page.png"),
    "chat": screenshot_path("tmp/fix-audit/06-chat.png", "tmp/product-pass/05-chat-buyer.png"),
    "manage": screenshot_path("tmp/screenshots/5_manage_page.png", "tmp/product-pass/02-manage-seller.png"),
}


def make_placeholder(name: str, size=(1280, 800)) -> Path:
    path = ASSETS / f"{name}_placeholder.png"
    im = Image.new("RGB", size, COLORS["bg"])
    draw = ImageDraw.Draw(im)
    for radius, color in [(460, COLORS["surface2"]), (310, COLORS["surface"]), (170, COLORS["teal"])]:
        box = (size[0] // 2 - radius, size[1] // 2 - radius, size[0] // 2 + radius, size[1] // 2 + radius)
        draw.ellipse(box, outline=color, width=4)
    draw.text((80, 80), "Radius project screenshot", fill=COLORS["text"])
    draw.text((80, 122), name.replace("_", " ").title(), fill=COLORS["muted"])
    im.save(path)
    return path


def prepare_image(name: str, src: Path | None, size=(1600, 900)) -> Path:
    if src is None:
        src = make_placeholder(name, size)
    target = ASSETS / f"{name}.png"
    im = Image.open(src).convert("RGB")
    im = ImageOps.contain(im, size)
    canvas = Image.new("RGB", size, COLORS["bg"])
    x = (size[0] - im.width) // 2
    y = (size[1] - im.height) // 2
    shadow = Image.new("RGBA", size, (0, 0, 0, 0))
    sh = Image.new("RGBA", (im.width, im.height), (0, 0, 0, 160)).filter(ImageFilter.GaussianBlur(18))
    shadow.alpha_composite(sh, (x + 12, y + 16))
    canvas = Image.alpha_composite(canvas.convert("RGBA"), shadow)
    canvas.alpha_composite(im.convert("RGBA"), (x, y))
    canvas.convert("RGB").save(target, quality=92)
    return target


def draw_architecture() -> Path:
    path = ASSETS / "architecture.png"
    fig, ax = plt.subplots(figsize=(13, 7), facecolor="#0d0a12")
    ax.set_facecolor("#0d0a12")
    ax.axis("off")
    boxes = [
        ("React Client\nTrust Radar, Browse, Sell,\nChat, Admin", 0.08, 0.58, "#00d9c0"),
        ("Express API\nAuth, Listings, Uploads,\nReports, Admin", 0.39, 0.58, "#3b82ff"),
        ("FastAPI ML Service\nFraud Score, pHash,\nPrice Suggestion", 0.70, 0.58, "#ff3d81"),
        ("Supabase Postgres\nUsers, Listings, Reports,\nML Predictions", 0.27, 0.15, "#ffc531"),
        ("Supabase Storage\nListing Photos", 0.58, 0.15, "#a855f7"),
    ]
    for text, x, y, color in boxes:
        rect = plt.Rectangle((x, y), 0.22, 0.22, facecolor="#1c1626", edgecolor=color, linewidth=2.4)
        ax.add_patch(rect)
        ax.text(x + 0.11, y + 0.11, text, ha="center", va="center", color="#f5f3ff", fontsize=12, weight="bold")
    arrows = [
        ((0.30, 0.69), (0.39, 0.69)),
        ((0.61, 0.69), (0.70, 0.69)),
        ((0.50, 0.58), (0.38, 0.37)),
        ((0.50, 0.58), (0.69, 0.37)),
    ]
    for start, end in arrows:
        ax.annotate("", xy=end, xytext=start, arrowprops=dict(arrowstyle="->", color="#a69abc", lw=2.2))
    ax.text(0.5, 0.93, "Radius System Architecture", ha="center", color="#f5f3ff", fontsize=22, weight="bold")
    ax.text(0.5, 0.88, "Marketplace UI + secure API + explainable ML trust service", ha="center", color="#a69abc", fontsize=13)
    fig.savefig(path, dpi=170, bbox_inches="tight", pad_inches=0.25)
    plt.close(fig)
    return path


def draw_fraud_pipeline() -> Path:
    path = ASSETS / "fraud_pipeline.png"
    fig, ax = plt.subplots(figsize=(13, 7), facecolor="#0d0a12")
    ax.set_facecolor("#0d0a12")
    ax.axis("off")
    steps = [
        ("Listing\nsubmitted", "#00d9c0"),
        ("Photo pHash\n+ duplicate check", "#3b82ff"),
        ("Price anomaly\nbaseline", "#ffc531"),
        ("Seller + text\nrisk signals", "#ff3d81"),
        ("Score 0-100\n+ explanations", "#a855f7"),
        ("Admin review\n+ ML labels", "#ff5370"),
    ]
    x0 = 0.05
    for i, (label, color) in enumerate(steps):
        x = x0 + i * 0.155
        circle = plt.Circle((x + 0.055, 0.55), 0.065, facecolor="#1c1626", edgecolor=color, linewidth=2.6)
        ax.add_patch(circle)
        ax.text(x + 0.055, 0.55, str(i + 1), ha="center", va="center", color="#f5f3ff", fontsize=16, weight="bold")
        ax.text(x + 0.055, 0.36, label, ha="center", va="center", color="#f5f3ff", fontsize=11, weight="bold")
        if i < len(steps) - 1:
            ax.annotate("", xy=(x + 0.145, 0.55), xytext=(x + 0.12, 0.55), arrowprops=dict(arrowstyle="->", color="#a69abc", lw=2.0))
    ax.text(0.5, 0.86, "Explainable Fraud Detection Pipeline", ha="center", color="#f5f3ff", fontsize=22, weight="bold")
    ax.text(0.5, 0.20, "Every prediction is saved with model version, component scores, signals, explanations, and feature snapshot hash.", ha="center", color="#a69abc", fontsize=12)
    fig.savefig(path, dpi=170, bbox_inches="tight", pad_inches=0.25)
    plt.close(fig)
    return path


def draw_dataset_chart() -> Path:
    path = ASSETS / "dataset_chart.png"
    labels = ["Price rows", "Server tests", "ML tests", "Categories", "Roles"]
    values = [520, 23, 9, 10, 3]
    colors = ["#00d9c0", "#3b82ff", "#ff3d81", "#ffc531", "#a855f7"]
    fig, ax = plt.subplots(figsize=(12, 6.5), facecolor="#0d0a12")
    ax.set_facecolor("#0d0a12")
    bars = ax.bar(labels, values, color=colors)
    ax.tick_params(colors="#f5f3ff", labelsize=11)
    ax.spines[:].set_visible(False)
    ax.grid(axis="y", color="#2f263d", alpha=0.8)
    ax.set_title("Project Evidence Snapshot", color="#f5f3ff", fontsize=22, weight="bold", pad=18)
    for bar, value in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, value + max(values) * 0.02, str(value), ha="center", color="#f5f3ff", weight="bold")
    fig.savefig(path, dpi=170, bbox_inches="tight", pad_inches=0.25)
    plt.close(fig)
    return path


def add_transition(slide):
    xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>'
    slide._element.append(parse_xml(xml))


def add_bg(slide, title=None):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("bg")
    # soft accent strips
    for left, top, width, height, color, alpha in [
        (0, 0, 13.33, 0.10, "teal", None),
        (0, 7.40, 13.33, 0.10, "pink", None),
    ]:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(left), PptInches(top), PptInches(width), PptInches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(color)
        shape.line.fill.background()
    if title:
        box = slide.shapes.add_textbox(PptInches(0.55), PptInches(0.30), PptInches(8.5), PptInches(0.55))
        p = box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = PptPt(24)
        p.font.bold = True
        p.font.color.rgb = rgb("text")


def add_text(slide, text, left, top, width, height, size=22, color="text", bold=False, align=None):
    box = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = PptPt(size)
        p.font.bold = bold
        p.font.color.rgb = rgb(color)
        if align:
            p.alignment = align
    return box


def add_card(slide, left, top, width, height, title, body, accent="teal"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("surface")
    shape.line.color.rgb = rgb(accent)
    shape.line.width = PptPt(1.5)
    add_text(slide, title, left + 0.18, top + 0.16, width - 0.36, 0.35, size=15, bold=True, color=accent)
    add_text(slide, body, left + 0.18, top + 0.56, width - 0.36, height - 0.68, size=12, color="text")


def add_image(slide, path: Path, left, top, width, height):
    slide.shapes.add_picture(str(path), PptInches(left), PptInches(top), width=PptInches(width), height=PptInches(height))


def build_pptx(assets: dict[str, Path]) -> Path:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    slides = []
    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s)
    add_text(s, "Radius", 0.65, 0.55, 5.2, 0.9, size=48, bold=True)
    add_text(s, "AI Assisted Hyperlocal Marketplace", 0.70, 1.45, 5.8, 0.45, size=20, color="teal", bold=True)
    add_text(s, "Buy and sell safely with people nearby.\nTrust Radar + explainable fraud detection + admin moderation.", 0.70, 2.15, 5.4, 1.0, size=18, color="text")
    add_card(s, 0.72, 5.45, 4.4, 0.9, "Presented by", f"{TEAM}\n{MEMBERS}", "blue")
    add_image(s, assets["hero"], 6.7, 0.45, 5.75, 6.45)

    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, "Problem")
    add_text(s, "Secondhand marketplaces are useful, but trust is weak.", 0.65, 1.0, 8.8, 0.5, size=28, bold=True)
    problems = [
        ("Fake listings", "Scammers post attractive products that do not exist.", "danger"),
        ("Copied photos", "Old or stolen product images mislead buyers.", "pink"),
        ("Suspicious prices", "Very low prices are used as bait.", "amber"),
        ("Off-platform pressure", "Users are pushed to WhatsApp, deposits, or token money.", "blue"),
    ]
    for i, (t, b, c) in enumerate(problems):
        add_card(s, 0.8 + (i % 2) * 5.8, 2.0 + (i // 2) * 1.85, 5.1, 1.28, t, b, c)

    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, "Solution")
    add_text(s, "Radius combines local discovery with AI trust screening.", 0.65, 0.95, 8.4, 0.5, size=27, bold=True)
    for i, item in enumerate(["Browse nearby products", "Sell with secure photo metadata", "Chat inside the platform", "Report suspicious listings", "Admin reviews AI flags", "ML log explains predictions"]):
        add_card(s, 0.75 + (i % 3) * 4.05, 1.8 + (i // 3) * 2.0, 3.55, 1.35, item, "Implemented in the current full-stack system.", ["teal", "blue", "pink", "amber", "danger", "teal"][i])

    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, "Architecture")
    add_image(s, assets["architecture"], 0.55, 0.92, 12.2, 5.9)

    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, "User Roles")
    add_card(s, 0.80, 1.45, 3.45, 4.55, "Guest", "Browse listings\nSearch and filter\nView listing details\nCannot sell/chat/report", "teal")
    add_card(s, 4.90, 1.45, 3.45, 4.55, "Registered User", "Create listings\nUpload photos\nChat with sellers\nReport suspicious items\nReview after completed trade", "blue")
    add_card(s, 9.00, 1.45, 3.45, 4.55, "Admin", "Review fraud queue\nResolve reports\nManage users\nInspect ML prediction log", "pink")

    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, "Trust Radar")
    add_image(s, assets["hero"], 0.65, 0.85, 6.0, 6.05)
    add_text(s, "The radar is a live discovery feature, not decoration.", 7.05, 1.05, 5.5, 0.55, size=25, bold=True)
    add_card(s, 7.15, 1.95, 4.95, 1.0, "Distance", "The center is the user; nodes are nearby products.", "teal")
    add_card(s, 7.15, 3.15, 4.95, 1.0, "Trust", "Node color shows verified, review, or flagged state.", "amber")
    add_card(s, 7.15, 4.35, 4.95, 1.0, "Interaction", "Selecting a node opens a mini listing preview card.", "blue")

    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, "Product Detail")
    add_image(s, assets["detail"], 0.65, 0.85, 6.0, 6.05)
    add_card(s, 7.10, 1.05, 4.95, 1.15, "Trust Scan", "Shows fraud score, signals, explanations, and visual risk bar.", "teal")
    add_card(s, 7.10, 2.55, 4.95, 1.15, "Price Radar", "Compares listing price against similar nearby products.", "amber")
    add_card(s, 7.10, 4.05, 4.95, 1.15, "Seller Safety", "Seller reputation and report workflow support safer decisions.", "pink")

    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, "Sell and Manage")
    add_image(s, assets["sell"], 0.65, 0.85, 5.65, 6.05)
    add_image(s, assets["manage"], 6.70, 1.18, 5.75, 3.60)
    add_text(s, "Sell flow triggers ML scoring. Manage tab focuses on seller operations.", 6.85, 5.18, 5.25, 0.8, size=20, bold=True)

    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, "Fraud Detection Engine")
    add_image(s, assets["fraud_pipeline"], 0.55, 0.92, 12.2, 5.9)

    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, "Admin Moderation")
    add_image(s, assets["admin"], 0.65, 0.85, 6.0, 6.05)
    add_card(s, 7.10, 1.05, 4.95, 1.05, "Fraud Queue", "Approve or remove high-risk listings.", "danger")
    add_card(s, 7.10, 2.35, 4.95, 1.05, "Reports", "Resolve user-submitted suspicious listing reports.", "amber")
    add_card(s, 7.10, 3.65, 4.95, 1.05, "Users", "Suspend or reactivate accounts.", "blue")
    add_card(s, 7.10, 4.95, 4.95, 1.05, "ML Log", "Audit model version, score, signals, and feature snapshot.", "teal")

    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, "Dataset and Evaluation")
    add_image(s, assets["dataset_chart"], 0.70, 1.05, 6.4, 4.95)
    add_card(s, 7.55, 1.12, 4.7, 1.05, "Price Dataset", "520+ market price rows across 10 product categories.", "teal")
    add_card(s, 7.55, 2.42, 4.7, 1.05, "Fraud Labels", "Admin decisions, reports, and manual CSV imports become labels.", "pink")
    add_card(s, 7.55, 3.72, 4.7, 1.05, "Metrics", "Precision, recall, F1 score, and false positive rate.", "amber")

    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, "Testing and Result")
    add_text(s, "Verified implementation", 0.72, 1.05, 5.6, 0.6, size=30, bold=True)
    add_card(s, 0.85, 2.0, 3.55, 1.3, "ML Service", "9 pytest tests passed", "teal")
    add_card(s, 4.90, 2.0, 3.55, 1.3, "Server", "23 Vitest tests passed", "blue")
    add_card(s, 8.95, 2.0, 3.55, 1.3, "Client", "Production Vite build passed", "pink")
    add_text(s, "The project satisfies the academic requirements and is structured as a portfolio-ready full-stack product.", 1.25, 4.45, 10.8, 0.8, size=23, align=PP_ALIGN.CENTER)

    s = prs.slides.add_slide(blank); slides.append(s); add_bg(s, "Conclusion")
    add_text(s, "Radius is more than a product selling system.", 0.95, 1.05, 8.2, 0.6, size=32, bold=True)
    add_text(s, "It combines role-based marketplace workflows with nearby discovery, real-time chat, admin moderation, and explainable AI/ML trust scoring.", 0.98, 2.05, 10.9, 1.0, size=24)
    add_card(s, 1.0, 4.05, 3.5, 1.2, "Local", "Products near the user.", "teal")
    add_card(s, 4.9, 4.05, 3.5, 1.2, "Safe", "Fraud signals are visible.", "blue")
    add_card(s, 8.8, 4.05, 3.5, 1.2, "Explainable", "Admins see reasons, not just scores.", "pink")
    add_text(s, "Thank you", 0.95, 6.35, 11.4, 0.45, size=22, color="teal", bold=True, align=PP_ALIGN.CENTER)

    for slide in slides:
        add_transition(slide)
    output = OUT / "Radius_AI_Hyperlocal_Marketplace_Presentation.pptx"
    prs.save(output)
    return output


def set_doc_styles(doc: Document):
    styles = doc.styles
    styles["Normal"].font.name = "Aptos"
    styles["Normal"].font.size = Pt(10.5)
    styles["Heading 1"].font.name = "Aptos Display"
    styles["Heading 1"].font.size = Pt(18)
    styles["Heading 1"].font.color.rgb = doc_rgb("teal")
    styles["Heading 2"].font.name = "Aptos Display"
    styles["Heading 2"].font.size = Pt(14)
    styles["Heading 2"].font.color.rgb = doc_rgb("blue")


def add_doc_picture(doc: Document, path: Path, width=6.2):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run()
    run.add_picture(str(path), width=Inches(width))


def add_bullets(doc: Document, items: list[str]):
    for item in items:
        doc.add_paragraph(item, style="List Bullet")


def build_report(assets: dict[str, Path]) -> Path:
    doc = Document()
    set_doc_styles(doc)
    section = doc.sections[0]
    section.top_margin = Inches(0.65)
    section.bottom_margin = Inches(0.65)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = title.add_run("Radius: AI Assisted Hyperlocal Marketplace")
    r.bold = True
    r.font.size = Pt(22)
    r.font.color.rgb = doc_rgb("teal")
    sub = doc.add_paragraph("Final Project Report")
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.runs[0].font.size = Pt(14)
    doc.add_paragraph(f"Team: {TEAM}\nMembers: {MEMBERS}\nCourse: {COURSE}\nInstructor: {INSTRUCTOR}\nDate: {DATE}")

    doc.add_heading("Abstract", level=1)
    doc.add_paragraph(
        "Radius is a full-stack hyperlocal secondhand marketplace that helps users buy and sell products near their location. "
        "The system supports guest, registered user, and admin roles. Registered users can create listings, upload images, chat, "
        "report suspicious listings, complete trades, and review users after completed trades. Admins can moderate reports, manage users, "
        "review suspicious listings, and inspect AI/ML prediction logs."
    )
    doc.add_paragraph(
        "The intelligent feature is an explainable trust and fraud detection engine. It analyzes product image similarity, price anomaly, "
        "seller history, listing text, brand consistency, reused descriptions, and off-platform contact signals. The system returns a score, "
        "decision, signal list, explanations, component scores, model version, and feature snapshot hash. The ML service also includes price "
        "suggestion with Random Forest regression and candidate fraud training with TF-IDF and Naive Bayes."
    )

    doc.add_heading("1. Introduction", level=1)
    doc.add_paragraph(
        "Secondhand marketplaces are useful, but they often suffer from fake listings, copied photos, suspicious prices, and sellers who push buyers "
        "to pay outside the platform. Radius addresses this by combining nearby product discovery with explainable trust screening."
    )
    add_doc_picture(doc, assets["hero"], 5.7)

    doc.add_heading("2. Objectives", level=1)
    add_bullets(doc, [
        "Build a role-based marketplace system with guest, registered user, and admin roles.",
        "Support product selling, browsing, searching, chat, reporting, trade, and review workflows.",
        "Use geofenced browsing so users can discover nearby products.",
        "Add an AI/ML based trust and fraud detection feature.",
        "Provide admin moderation and ML prediction audit logs.",
        "Create a polished portfolio-quality project with tests and documentation.",
    ])

    doc.add_heading("3. Requirement Mapping", level=1)
    table = doc.add_table(rows=1, cols=2)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    hdr = table.rows[0].cells
    hdr[0].text = "Requirement"
    hdr[1].text = "Implemented in Radius"
    rows = [
        ("Different user roles", "Guest, registered user, admin"),
        ("Product selling system", "Listing create/manage, image upload, search, detail page"),
        ("AI/ML feature", "Trust scoring, duplicate image detection, price anomaly, price suggestion, candidate fraud training"),
        ("Dataset", "Market price CSV, manual label CSV, admin/report labels"),
        ("Python ML", "FastAPI, scikit-learn, pandas, Pillow, imagehash"),
        ("Frameworks", "React, Express, Supabase, FastAPI, Socket.io"),
    ]
    for left, right in rows:
        cells = table.add_row().cells
        cells[0].text = left
        cells[1].text = right
        for cell in cells:
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.TOP

    doc.add_heading("4. Literature Review", level=1)
    doc.add_paragraph(
        "Fraud detection literature shows that marketplace fraud can be detected using seller behavior, listing behavior, image similarity, and transaction signals. "
        "The project follows research guidance by storing features, labels, predictions, and evaluation metrics rather than using accuracy alone."
    )
    add_bullets(doc, [
        "Mutemi and Bacao (2023) show marketplace fraud can be detected using machine learning features from digital marketplaces.",
        "Ali et al. (2022) explain fraud detection metrics such as precision, recall, F1, and false positive rate.",
        "Chatrath, Batra, and Chaba (2022) motivate image-based checking for secondhand product trust.",
        "Amazon Science Fraud Dataset Benchmark motivates standardized datasets, splits, and evaluation practices.",
    ])

    doc.add_heading("5. System Architecture", level=1)
    doc.add_paragraph("Radius has three main layers: React frontend, Express backend API, and Python FastAPI ML service. Supabase Postgres stores persistent data and PostGIS supports location search.")
    add_doc_picture(doc, assets["architecture"], 6.5)

    doc.add_heading("6. User Roles", level=1)
    add_bullets(doc, [
        "Guest: browse, search, filter, and view listing details.",
        "Registered user: create listings, upload photos, chat, report suspicious listings, manage own listings, and review after completed trades.",
        "Admin: view analytics, review fraud queue, resolve reports, manage users, and inspect ML prediction logs.",
    ])

    doc.add_heading("7. Main Features", level=1)
    doc.add_heading("7.1 Trust Radar", level=2)
    doc.add_paragraph("Trust Radar is an interactive product discovery view. The center represents user location, and nearby products appear as image nodes with trust colors.")
    doc.add_heading("7.2 Listing and Sell Flow", level=2)
    doc.add_paragraph("Registered users can submit product information and photos. The listing is scored by the ML service before being stored.")
    add_doc_picture(doc, assets["sell"], 5.2)
    doc.add_heading("7.3 Listing Detail", level=2)
    doc.add_paragraph("The listing detail page shows product information, seller information, trust score, price comparison, report action, and chat action.")
    add_doc_picture(doc, assets["detail"], 5.2)
    doc.add_heading("7.4 Chat, Trade, and Review", level=2)
    doc.add_paragraph("Socket.io powers real-time chat. Reviews are allowed only after completed trades, which reduces fake reputation.")
    add_doc_picture(doc, assets["chat"], 5.2)
    doc.add_heading("7.5 Admin Dashboard", level=2)
    doc.add_paragraph("Admin can review fraud queue, resolve reports, manage users, and inspect ML prediction logs.")
    add_doc_picture(doc, assets["admin"], 5.2)

    doc.add_heading("8. AI/ML Feature", level=1)
    doc.add_paragraph("The fraud engine uses explainable signals. Each risky signal adds points and trusted seller history can reduce points. The final score is clamped from 0 to 100.")
    add_doc_picture(doc, assets["fraud_pipeline"], 6.5)
    add_bullets(doc, [
        "Duplicate or visually similar image using pHash and Hamming distance.",
        "Price anomaly using category-condition baseline.",
        "New seller high-value listing and seller prior flags.",
        "Urgent language, payment pressure, off-platform contact, and prohibited item words.",
        "Brand/title mismatch and reused title-description.",
        "Feature snapshot, component scores, model version, and explanations saved for audit.",
    ])

    doc.add_heading("9. Dataset and Evaluation", level=1)
    doc.add_paragraph("The current dataset includes market price rows across product categories. Admin decisions, reports, and manual CSV imports can produce fraud labels for future training.")
    add_doc_picture(doc, assets["dataset_chart"], 6.0)
    add_bullets(doc, [
        "Price suggestion uses Random Forest regression.",
        "Candidate fraud model uses TF-IDF and Multinomial Naive Bayes.",
        "Evaluation metrics include precision, recall, F1, and false positive rate.",
        "The online decision is explainable and human-reviewed rather than fully automatic removal.",
    ])

    doc.add_heading("10. Security", level=1)
    add_bullets(doc, [
        "Passwords are hashed with bcrypt.",
        "JWT is used for authentication.",
        "Admin routes require admin role.",
        "Listing mutation checks owner or admin permission.",
        "Supabase service role key stays server-side only.",
        "Express uses validation, CORS, rate limiting, and Helmet.",
    ])

    doc.add_heading("11. Testing and Results", level=1)
    add_bullets(doc, [
        "ML service tests: 9 passed.",
        "Server tests: 23 passed.",
        "Client production build: passed.",
    ])

    doc.add_heading("12. Limitations", level=1)
    add_bullets(doc, [
        "The dataset is demonstration-sized and should be expanded before production.",
        "The online fraud engine is an explainable hybrid scoring system, not a large production-trained deep learning model.",
        "Payment gateway and identity verification are not included in the academic scope.",
    ])

    doc.add_heading("13. Future Work", level=1)
    add_bullets(doc, [
        "Collect larger labeled marketplace datasets ethically.",
        "Train and deploy a stronger supervised fraud classifier.",
        "Add graph-based fraud detection for seller-buyer networks.",
        "Add escrow payment, identity verification, and mobile app support.",
    ])

    doc.add_heading("14. Conclusion", level=1)
    doc.add_paragraph(
        "Radius satisfies the project requirements by implementing role-based access, a product selling system, a Python AI/ML service, dataset support, "
        "admin moderation, and a polished user interface. Its main contribution is combining nearby product discovery with explainable fraud detection."
    )

    doc.add_heading("References", level=1)
    refs = [
        "Mutemi, A., and Bacao, F. (2023). A numeric-based machine learning design for detecting organized retail fraud in digital marketplaces. Scientific Reports. https://doi.org/10.1038/s41598-023-38304-5",
        "Ali, A. et al. (2022). Financial Fraud Detection Based on Machine Learning: A Systematic Literature Review. Applied Sciences. https://doi.org/10.3390/app12199637",
        "Chatrath, S. K., Batra, G. S., and Chaba, Y. (2022). Handling consumer vulnerability in e-commerce product images using machine learning. Heliyon. https://doi.org/10.1016/j.heliyon.2022.e10743",
        "Amazon Science. Fraud Dataset Benchmark. https://github.com/amazon-science/fraud-dataset-benchmark",
    ]
    for ref in refs:
        doc.add_paragraph(ref, style="List Number")

    output = OUT / "Radius_AI_Hyperlocal_Marketplace_Report.docx"
    doc.save(output)
    return output


def main():
    prepared = {name: prepare_image(name, path) for name, path in SCREENSHOTS.items()}
    prepared["architecture"] = draw_architecture()
    prepared["fraud_pipeline"] = draw_fraud_pipeline()
    prepared["dataset_chart"] = draw_dataset_chart()
    pptx = build_pptx(prepared)
    report = build_report(prepared)
    print(pptx)
    print(report)


if __name__ == "__main__":
    main()
